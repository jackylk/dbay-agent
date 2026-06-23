"""Merge new-slides-white.pptx (prepend) + memory-evolution.pptx using OOXML copy."""
import copy
from pptx import Presentation
from pptx.util import Emu

def append_slides(target_prs, source_path):
    src = Presentation(source_path)
    # Use last layout available
    layout = target_prs.slide_layouts[len(target_prs.slide_layouts) - 1]

    for src_slide in src.slides:
        new_slide = target_prs.slides.add_slide(layout)

        # Remove default placeholder shapes
        for ph in list(new_slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Copy all elements from source slide
        for el in src_slide.shapes._spTree:
            new_slide.shapes._spTree.append(copy.deepcopy(el))

        # Copy images
        for rel in src_slide.part.rels.values():
            if "image" in rel.reltype:
                new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

def main():
    new_slides = 'new-slides-white.pptx'
    existing = '/Users/jacky/code/lakeon/docs/reports/memory-evolution.pptx'
    output = '/Users/jacky/code/lakeon/docs/reports/memory-evolution-merged.pptx'

    prs = Presentation(new_slides)
    append_slides(prs, existing)
    prs.save(output)
    print(f'Saved: {output}')

main()
