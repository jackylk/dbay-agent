"""Document parsing: PDF, DOCX, DOC, XLSX/XLS, PPTX, EPUB, HTML, Markdown, Text."""
import logging
from typing import List, Dict, Tuple

logger = logging.getLogger(__name__)

# Max stream size for DOC OLE2 files (50MB)
_MAX_STREAM_SIZE = 50 * 1024 * 1024
_MAX_CCP_TEXT = 10_000_000


def parse_document(file_path: str, format: str) -> Tuple[str, List[Dict]]:
    """Parse a document and return (markdown_text, page_metadata).

    page_metadata: [{page: int, char_start: int, char_end: int}, ...]
    For non-PDF formats, page_metadata is empty.
    """
    format = format.upper()
    if format == "PDF":
        return _parse_pdf(file_path)
    elif format == "DOCX":
        return _parse_docx(file_path), []
    elif format == "DOC":
        return _parse_doc(file_path), []
    elif format in ("XLSX", "XLS", "XLSM"):
        return _parse_excel(file_path, format), []
    elif format == "PPTX":
        return _parse_pptx(file_path), []
    elif format == "EPUB":
        return _parse_epub(file_path), []
    elif format == "HTML":
        return _parse_html(file_path), []
    elif format in ("MARKDOWN", "MD", "TEXT", "TXT"):
        return _parse_markdown(file_path), []
    else:
        raise ValueError(f"Unsupported format: {format}")


def parse_document_to_fulltext(file_path: str, format: str) -> str:
    """Parse document and return just the markdown text (no page metadata)."""
    text, _ = parse_document(file_path, format)
    return text


# ── PDF ──────────────────────────────────────────────────────────────────

def _parse_pdf(file_path: str) -> Tuple[str, List[Dict]]:
    import pymupdf4llm
    import pymupdf

    doc = pymupdf.open(file_path)
    doc.close()

    pages = pymupdf4llm.to_markdown(file_path, page_chunks=True)

    page_metadata = []
    parts = []
    current_offset = 0

    for i, page_data in enumerate(pages):
        page_num = page_data.get("metadata", {}).get("page", 0)
        text = page_data.get("text", "")
        if not text:
            continue
        char_start = current_offset
        parts.append(text)
        current_offset += len(text)
        if i < len(pages) - 1:
            parts.append("\n\n")
            current_offset += 2
        page_metadata.append({
            "page": page_num + 1,
            "char_start": char_start,
            "char_end": current_offset,
        })

    markdown = "".join(parts)
    return markdown, page_metadata


# ── DOCX (upgraded: preserves paragraph/table order + formatting) ────────

def _parse_docx(file_path: str) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph

    doc = Document(file_path)
    markdown_parts = []

    # Map XML table elements for O(1) lookup
    table_by_element = {table._tbl: table for table in doc.tables}

    # Walk document body in order (preserves table positions)
    for child in doc.element.body:
        if child.tag == qn("w:p"):
            paragraph = Paragraph(child, doc)
            if not paragraph.text.strip():
                continue
            style_name = paragraph.style.name if paragraph.style else "Normal"
            if style_name.startswith("Heading"):
                level = _extract_heading_level(style_name)
                markdown_parts.append(f"{'#' * level} {paragraph.text}")
            else:
                text = _convert_formatted_text(paragraph)
                markdown_parts.append(text)
        elif child.tag == qn("w:tbl"):
            if child in table_by_element:
                markdown_parts.append(_convert_word_table(table_by_element[child]))

    return "\n\n".join(markdown_parts)


def _extract_heading_level(style_name: str) -> int:
    try:
        for part in style_name.split():
            if part.isdigit():
                return min(int(part), 6)
    except Exception:
        pass
    return 1


def _convert_formatted_text(paragraph) -> str:
    text_parts = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        if run.bold:
            text = f"**{text}**"
        if run.italic:
            text = f"*{text}*"
        text_parts.append(text)
    return "".join(text_parts)


def _convert_word_table(table) -> str:
    if not table.rows:
        return ""
    rows = []
    for row in table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        rows.append(row_data)
    return _format_table_to_markdown(rows)


# ── DOC (OLE2 legacy Word 97-2003) ──────────────────────────────────────

def _parse_doc(file_path: str) -> str:
    import struct
    import olefile

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception as e:
        logger.warning(f"Failed to open .doc as OLE file: {e}")
        return _doc_fallback_extract(file_path)

    try:
        return _doc_extract_from_ole(ole, struct)
    except Exception as e:
        logger.warning(f"Structured OLE extraction failed, using fallback: {e}")
        return _doc_fallback_extract(file_path)
    finally:
        ole.close()


def _doc_read_stream(ole, name: str) -> bytes:
    stream = ole.openstream(name)
    data = stream.read(_MAX_STREAM_SIZE + 1)
    if len(data) > _MAX_STREAM_SIZE:
        raise ValueError(f"OLE stream '{name}' exceeds {_MAX_STREAM_SIZE} bytes")
    return data


def _doc_extract_from_ole(ole, struct) -> str:
    if not ole.exists("WordDocument"):
        raise ValueError("No WordDocument stream found")

    word_doc = _doc_read_stream(ole, "WordDocument")
    if len(word_doc) < 0x01A8:
        raise ValueError(f"WordDocument stream too small ({len(word_doc)} bytes)")

    nfib = struct.unpack_from("<H", word_doc, 0x0002)[0]
    if nfib < 0x00C1:
        raise ValueError(f"Unsupported Word version (nFib=0x{nfib:04X})")

    flags = struct.unpack_from("<H", word_doc, 0x000A)[0]
    table_stream = "1Table" if (flags & 0x0200) else "0Table"

    ccp_text = struct.unpack_from("<i", word_doc, 0x004C)[0]
    if ccp_text <= 0:
        raise ValueError("ccpText is zero or negative")
    ccp_text = min(ccp_text, _MAX_CCP_TEXT)

    if not ole.exists(table_stream):
        raise ValueError(f"Table stream '{table_stream}' not found")
    table_data = _doc_read_stream(ole, table_stream)

    fc_clx = struct.unpack_from("<i", word_doc, 0x01A2)[0]
    lcb_clx = struct.unpack_from("<i", word_doc, 0x01A6)[0]

    if fc_clx <= 0 or lcb_clx <= 0 or fc_clx + lcb_clx > len(table_data):
        return _doc_simple_extract(word_doc, ccp_text, struct)

    return _doc_extract_via_clx(word_doc, table_data, fc_clx, lcb_clx, ccp_text, struct)


def _doc_simple_extract(word_doc: bytes, ccp_text: int, struct) -> str:
    text_start = 0x0800
    if text_start >= len(word_doc):
        raise ValueError("WordDocument stream too small for text extraction")

    if ccp_text * 2 + text_start <= len(word_doc):
        end = text_start + ccp_text * 2
        text = word_doc[text_start:end].decode("utf-16-le", errors="replace")
        if sum(1 for c in text[:200] if c.isprintable() or c in "\n\r\t") > len(text[:200]) * 0.5:
            return _doc_clean_text(text)

    end = min(text_start + ccp_text, len(word_doc))
    return _doc_clean_text(word_doc[text_start:end].decode("cp1252", errors="replace"))


def _doc_extract_via_clx(word_doc, table_data, fc_clx, lcb_clx, ccp_text, struct):
    clx = table_data[fc_clx:fc_clx + lcb_clx]
    pos = 0
    text_parts = []
    chars_extracted = 0

    while pos < len(clx) and clx[pos] == 0x01:
        if pos + 3 > len(clx):
            break
        cb = struct.unpack_from("<H", clx, pos + 1)[0]
        advance = 3 + cb
        if advance <= 0:
            break
        pos += advance

    if pos >= len(clx) or clx[pos] != 0x02:
        return _doc_simple_extract(word_doc, ccp_text, struct)

    pos += 1
    if pos + 4 > len(clx):
        return _doc_simple_extract(word_doc, ccp_text, struct)

    lcb_pcd = struct.unpack_from("<I", clx, pos)[0]
    pos += 4
    pcd_end = pos + lcb_pcd

    if pcd_end > len(clx):
        return _doc_simple_extract(word_doc, ccp_text, struct)

    n_pieces = (lcb_pcd - 4) // 12
    if n_pieces <= 0:
        return _doc_simple_extract(word_doc, ccp_text, struct)

    cps = []
    for i in range(n_pieces + 1):
        offset = pos + i * 4
        if offset + 4 > len(clx):
            break
        cps.append(struct.unpack_from("<I", clx, offset)[0])

    pcd_array_start = pos + (n_pieces + 1) * 4

    for i in range(min(n_pieces, len(cps) - 1)):
        if chars_extracted >= ccp_text:
            break
        pcd_offset = pcd_array_start + i * 8
        if pcd_offset + 8 > len(clx):
            break

        fc_value = struct.unpack_from("<I", clx, pcd_offset + 2)[0]
        piece_cp_start = cps[i]
        piece_cp_end = cps[i + 1]
        piece_char_count = piece_cp_end - piece_cp_start

        is_compressed = bool(fc_value & 0x40000000)
        fc_real = fc_value & 0x3FFFFFFF

        if is_compressed:
            byte_offset = fc_real // 2
            byte_end = byte_offset + piece_char_count
            if byte_end <= len(word_doc):
                text_parts.append(word_doc[byte_offset:byte_end].decode("cp1252", errors="replace"))
        else:
            byte_offset = fc_real
            byte_end = byte_offset + piece_char_count * 2
            if byte_end <= len(word_doc):
                text_parts.append(word_doc[byte_offset:byte_end].decode("utf-16-le", errors="replace"))

        chars_extracted += piece_char_count

    result = _doc_clean_text("".join(text_parts))
    if not result.strip():
        return _doc_simple_extract(word_doc, ccp_text, struct)
    return result


def _doc_clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\x07", "\t").replace("\x0B", "\n").replace("\x0C", "\n\n")
    return text


def _doc_fallback_extract(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw = f.read(_MAX_STREAM_SIZE)

    try:
        decoded = raw.decode("utf-16-le", errors="ignore")
        lines = []
        current = []
        for ch in decoded:
            if ch.isprintable() or ch in "\n\t":
                current.append(ch)
            else:
                if len(current) > 3:
                    lines.append("".join(current))
                current = []
        if current and len(current) > 3:
            lines.append("".join(current))
        text = "\n".join(lines)
        if len(text) > 50:
            return text
    except Exception:
        pass

    text = raw.decode("cp1252", errors="replace")
    lines = []
    current = []
    for ch in text:
        if ch.isprintable() or ch in "\n\t":
            current.append(ch)
        else:
            if len(current) > 3:
                lines.append("".join(current))
            current = []
    if current and len(current) > 3:
        lines.append("".join(current))
    return "\n".join(lines)


# ── XLSX / XLS / XLSM ───────────────────────────────────────────────────

MAX_ROWS_PER_SHEET = 1000

def _parse_excel(file_path: str, format: str) -> str:
    if format == "XLS":
        return _parse_xls(file_path)
    else:
        return _parse_xlsx(file_path)


def _parse_xlsx(file_path: str) -> str:
    import openpyxl
    from pathlib import Path

    wb = openpyxl.load_workbook(file_path, data_only=True)
    parts = [f"# {Path(file_path).stem}", f"**Sheets:** {len(wb.sheetnames)}"]

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_parts = [f"## Sheet: {sheet_name}"]

        max_row = sheet.max_row or 0
        max_col = sheet.max_column or 0
        if max_row == 0 or max_col == 0:
            sheet_parts.append("*Empty sheet*")
            parts.append("\n\n".join(sheet_parts))
            continue

        sheet_parts.append(f"**Dimensions:** {max_row} rows × {max_col} columns")

        rows_to_process = min(max_row, MAX_ROWS_PER_SHEET)
        rows = []
        for row in sheet.iter_rows(min_row=1, max_row=rows_to_process, values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])

        if rows:
            sheet_parts.append(_format_table_to_markdown(rows))

        if max_row > MAX_ROWS_PER_SHEET:
            sheet_parts.append(f"\n*... {max_row - MAX_ROWS_PER_SHEET} more rows truncated ...*")

        parts.append("\n\n".join(sheet_parts))

    return "\n\n".join(parts)


def _parse_xls(file_path: str) -> str:
    import xlrd
    from pathlib import Path

    wb = xlrd.open_workbook(file_path, formatting_info=True, on_demand=True)
    try:
        parts = [f"# {Path(file_path).stem}", f"**Sheets:** {wb.nsheets}"]

        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            sheet_parts = [f"## Sheet: {sheet.name}"]

            if sheet.nrows == 0 or sheet.ncols == 0:
                sheet_parts.append("*Empty sheet*")
                parts.append("\n\n".join(sheet_parts))
                continue

            sheet_parts.append(f"**Dimensions:** {sheet.nrows} rows × {sheet.ncols} columns")

            rows_to_process = min(sheet.nrows, MAX_ROWS_PER_SHEET)
            rows = []
            for row_idx in range(rows_to_process):
                row_data = []
                for col_idx in range(sheet.ncols):
                    row_data.append(_format_xls_cell(sheet.cell(row_idx, col_idx), wb, xlrd))
                rows.append(row_data)

            if rows:
                sheet_parts.append(_format_table_to_markdown(rows))

            if sheet.nrows > MAX_ROWS_PER_SHEET:
                sheet_parts.append(f"\n*... {sheet.nrows - MAX_ROWS_PER_SHEET} more rows truncated ...*")

            parts.append("\n\n".join(sheet_parts))

        return "\n\n".join(parts)
    finally:
        wb.release_resources()


def _format_xls_cell(cell, wb, xlrd) -> str:
    if cell.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
        return ""
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            dt = xlrd.xldate_as_tuple(cell.value, wb.datemode)
            if dt[3] or dt[4] or dt[5]:
                return f"{dt[0]:04d}-{dt[1]:02d}-{dt[2]:02d} {dt[3]:02d}:{dt[4]:02d}:{dt[5]:02d}"
            return f"{dt[0]:04d}-{dt[1]:02d}-{dt[2]:02d}"
        except Exception:
            return str(cell.value)
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if cell.value else "FALSE"
    if cell.ctype == xlrd.XL_CELL_ERROR:
        error_map = {0x00: "#NULL!", 0x07: "#DIV/0!", 0x0F: "#VALUE!",
                     0x17: "#REF!", 0x1D: "#NAME?", 0x24: "#NUM!", 0x2A: "#N/A"}
        return error_map.get(cell.value, f"#ERR({cell.value})")
    if cell.ctype == xlrd.XL_CELL_NUMBER:
        if cell.value == int(cell.value):
            return str(int(cell.value))
        return str(cell.value)
    return str(cell.value) if cell.value is not None else ""


# ── PPTX ─────────────────────────────────────────────────────────────────

def _parse_pptx(file_path: str) -> str:
    import pptx
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = pptx.Presentation(file_path)
    parts = []
    slide_count = len(prs.slides)

    for idx, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {idx}/{slide_count}"]

        # Extract title
        title = ""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    title = shape.text.strip()
                    break
        if title:
            slide_parts.append(f"### {title}")

        # Extract content (non-title shapes)
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    continue
            if hasattr(shape, "text") and shape.text.strip():
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    slide_parts.append(_format_table_to_markdown(rows))
                else:
                    slide_parts.append(shape.text.strip())

        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"**Notes:** {notes}")

        parts.append("\n\n".join(slide_parts))

    return "\n\n---\n\n".join(parts)


# ── EPUB ─────────────────────────────────────────────────────────────────

def _parse_epub(file_path: str) -> str:
    try:
        return _parse_epub_ebooklib(file_path)
    except Exception as e:
        logger.warning(f"ebooklib failed, trying zipfile fallback: {e}")
        return _parse_epub_zipfile(file_path)


def _parse_epub_ebooklib(file_path: str) -> str:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(file_path, options={"ignore_ncx": True})
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
            level = int(tag.name[1])
            tag.replace_with(f"\n{'#' * level} {tag.get_text().strip()}\n")
        text = soup.get_text(separator="\n").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _parse_epub_zipfile(file_path: str) -> str:
    import zipfile
    import re
    import html as html_module

    parts = []
    with zipfile.ZipFile(file_path, "r") as zf:
        html_files = sorted([n for n in zf.namelist()
                            if n.endswith((".html", ".xhtml", ".htm"))])
        for name in html_files:
            content = zf.read(name).decode("utf-8", errors="replace")
            # Strip script/style
            content = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL)
            content = re.sub(r"<style[^>]*>.*?</style>", "", content, flags=re.DOTALL)
            # Convert headings
            for level in range(1, 7):
                content = re.sub(
                    rf"<h{level}[^>]*>(.*?)</h{level}>",
                    rf"\n{'#' * level} \1\n", content, flags=re.DOTALL | re.IGNORECASE)
            # Convert formatting
            content = re.sub(r"<(strong|b)[^>]*>(.*?)</\1>", r"**\2**", content, flags=re.DOTALL)
            content = re.sub(r"<(em|i)[^>]*>(.*?)</\1>", r"*\2*", content, flags=re.DOTALL)
            content = re.sub(r"<p[^>]*>(.*?)</p>", r"\1\n\n", content, flags=re.DOTALL)
            content = re.sub(r"<br\s*/?>", "\n", content)
            # Remove remaining tags
            content = re.sub(r"<[^>]+>", "", content)
            content = html_module.unescape(content)
            # Normalize whitespace
            content = re.sub(r"\n{3,}", "\n\n", content).strip()
            if content:
                parts.append(content)
    return "\n\n".join(parts)


# ── HTML ─────────────────────────────────────────────────────────────────

def _parse_html(file_path: str) -> str:
    """Parse local HTML file to markdown."""
    for encoding in ["utf-8", "utf-8-sig", "gbk", "gb2312", "latin1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                html_content = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        with open(file_path, "rb") as f:
            html_content = f.read().decode("utf-8", errors="replace")

    try:
        # Try readabilipy for main content extraction (removes ads/nav)
        from readabilipy import simple_json_from_html_string
        article = simple_json_from_html_string(html_content, use_readability=True)
        main_html = article.get("content", html_content)
    except Exception:
        main_html = html_content

    try:
        # Convert to markdown
        from markdownify import markdownify
        return markdownify(main_html, heading_style="ATX", strip=["script", "style"])
    except ImportError:
        # Fallback: basic HTML to text
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(main_html, "html.parser")
        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
            level = int(tag.name[1])
            tag.replace_with(f"\n{'#' * level} {tag.get_text().strip()}\n")
        return soup.get_text(separator="\n").strip()


# ── Markdown / Text ──────────────────────────────────────────────────────

def _parse_markdown(file_path: str) -> str:
    for encoding in ["utf-8", "gbk", "gb2312", "latin1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    with open(file_path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


# ── Shared utilities ─────────────────────────────────────────────────────

def _format_table_to_markdown(rows: list, has_header: bool = True) -> str:
    """Convert rows to markdown table."""
    if not rows:
        return ""

    # Calculate column widths
    num_cols = max(len(row) for row in rows)
    col_widths = [3] * num_cols
    for row in rows:
        for i, cell in enumerate(row):
            col_widths[i] = max(col_widths[i], len(str(cell)))

    lines = []
    for idx, row in enumerate(rows):
        # Pad row to num_cols
        padded = list(row) + [""] * (num_cols - len(row))
        cells = [str(c).ljust(col_widths[i]) for i, c in enumerate(padded)]
        lines.append("| " + " | ".join(cells) + " |")
        if idx == 0 and has_header:
            sep = ["-" * col_widths[i] for i in range(num_cols)]
            lines.append("| " + " | ".join(sep) + " |")

    return "\n".join(lines)
